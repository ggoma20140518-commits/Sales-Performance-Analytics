import streamlit as st
import pandas as pd
import plotly.express as px
import plotly.graph_objects as go
import re
import io
import datetime
import matplotlib
matplotlib.use("Agg")
import matplotlib.pyplot as plt
import matplotlib.font_manager as fm
import tempfile
import os
from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from pptx.enum.chart import XL_CHART_TYPE, XL_LEGEND_POSITION

# 한글 폰트 설정 (matplotlib) - Windows / Linux 대응
_font_path = None
_font_path_bold = None
_font_candidates = [
    r"C:\Windows\Fonts\malgun.ttf",          # Windows
    "/usr/share/fonts/truetype/nanum/NanumGothic.ttf",  # Linux (apt)
    os.path.join(os.path.dirname(__file__), "fonts", "NanumGothic.ttf"),  # 번들
]
_font_bold_candidates = [
    r"C:\Windows\Fonts\malgunbd.ttf",
    "/usr/share/fonts/truetype/nanum/NanumGothicBold.ttf",
    os.path.join(os.path.dirname(__file__), "fonts", "NanumGothicBold.ttf"),
]
for fp in _font_candidates:
    if os.path.exists(fp):
        _font_path = fp
        fm.fontManager.addfont(fp)
        plt.rcParams["font.family"] = fm.FontProperties(fname=fp).get_name()
        break
for fp in _font_bold_candidates:
    if os.path.exists(fp):
        _font_path_bold = fp
        break
plt.rcParams["axes.unicode_minus"] = False

st.set_page_config(page_title="COSMAX R&I 매출분석", layout="wide")

# ─── 커스텀 CSS (COSMAX PI 스타일) ───
st.markdown("""
<style>
@import url('https://fonts.googleapis.com/css2?family=Noto+Sans+KR:wght@300;400;500;700;900&display=swap');

/* 전체 폰트 */
html, body, [class*="css"], .stMarkdown, .stMetric, .stTabs, .stDataFrame {
    font-family: 'Noto Sans KR', 'Malgun Gothic', sans-serif !important;
}

/* 메인 배경 */
.stApp {
    background-color: #F9FAFB;
}

/* 사이드바 */
section[data-testid="stSidebar"] {
    background: linear-gradient(180deg, #111827 0%, #1F2937 100%);
}
section[data-testid="stSidebar"] * {
    color: #E5E7EB !important;
}
section[data-testid="stSidebar"] .stSelectbox label,
section[data-testid="stSidebar"] .stRadio label,
section[data-testid="stSidebar"] .stSlider label,
section[data-testid="stSidebar"] h2 {
    color: #F9FAFB !important;
    font-weight: 700 !important;
}
section[data-testid="stSidebar"] .stSelectbox [data-baseweb="select"],
section[data-testid="stSidebar"] .stRadio [data-baseweb="radio"],
section[data-testid="stSidebar"] input {
    background-color: #FFFFFF;
    border-color: #D1D5DB;
}
section[data-testid="stSidebar"] .stSelectbox [data-baseweb="select"] * {
    color: #1F2937 !important;
}
section[data-testid="stSidebar"] .stRadio [data-baseweb="radio"] {
    background-color: transparent;
    border-color: transparent;
}
section[data-testid="stSidebar"] .stRadio [role="radiogroup"] label span,
section[data-testid="stSidebar"] .stRadio [role="radiogroup"] label p {
    color: #F9FAFB !important;
}
section[data-testid="stSidebar"] hr {
    border-color: #374151;
}
section[data-testid="stSidebar"] button[kind="primary"] {
    background-color: #E61E3D !important;
    border: none !important;
    font-weight: 700 !important;
}
section[data-testid="stSidebar"] button[kind="primary"]:hover {
    background-color: #C41830 !important;
}

/* 헤더 영역 */
.main-header {
    background: linear-gradient(135deg, #111827 0%, #1F2937 60%, #E61E3D 100%);
    padding: 1.5rem 2rem;
    border-radius: 12px;
    margin-bottom: 1.5rem;
    position: relative;
    overflow: hidden;
}
.main-header::before {
    content: '';
    position: absolute;
    top: 0; right: 0;
    width: 40%;
    height: 100%;
    background: linear-gradient(135deg, transparent 30%, rgba(230,30,61,0.15) 100%);
}
.main-header h1 {
    color: #FFFFFF !important;
    font-size: 1.8rem !important;
    font-weight: 900 !important;
    margin: 0 !important;
    letter-spacing: -0.5px;
    display: flex;
    align-items: center;
    gap: 0.5rem;
}
.main-header h1 img {
    height: 60px;
    object-fit: contain;
}
.main-header p {
    color: #9CA3AF;
    font-size: 0.85rem;
    margin: 0.3rem 0 0 0;
}
.accent-bar {
    width: 60px;
    height: 4px;
    background: #E61E3D;
    border-radius: 2px;
    margin-top: 0.6rem;
}

/* KPI 카드 */
div[data-testid="stMetric"] {
    background: #FFFFFF;
    border: 1px solid #E5E7EB;
    border-radius: 10px;
    padding: 1rem 1.2rem;
    box-shadow: 0 1px 3px rgba(0,0,0,0.06);
    border-left: 4px solid #E61E3D;
}
div[data-testid="stMetric"] label {
    color: #6B7280 !important;
    font-size: 0.8rem !important;
    font-weight: 500 !important;
    text-transform: uppercase;
    letter-spacing: 0.5px;
}
div[data-testid="stMetric"] [data-testid="stMetricValue"] {
    color: #111827 !important;
    font-weight: 900 !important;
    font-size: 1.6rem !important;
}

/* 탭 스타일 */
.stTabs [data-baseweb="tab-list"] {
    background: #FFFFFF;
    border-radius: 10px;
    padding: 4px;
    gap: 4px;
    border: 1px solid #E5E7EB;
}
.stTabs [data-baseweb="tab"] {
    border-radius: 8px;
    font-weight: 600;
    color: #6B7280;
    padding: 0.5rem 1rem;
}
.stTabs [aria-selected="true"] {
    background-color: #111827 !important;
    color: #FFFFFF !important;
    border-radius: 8px;
}
.stTabs [data-baseweb="tab-panel"] {
    background: #FFFFFF;
    border-radius: 0 0 10px 10px;
    padding: 1.5rem;
    border: 1px solid #E5E7EB;
    border-top: none;
}

/* 서브헤더 */
h3, .stSubheader {
    color: #1F2937 !important;
    font-weight: 700 !important;
    border-left: 4px solid #E61E3D;
    padding-left: 0.8rem !important;
}

/* 데이터프레임 */
.stDataFrame {
    border-radius: 8px;
    overflow: hidden;
}

/* 구분선 */
hr {
    border: none;
    border-top: 1px solid #E5E7EB;
    margin: 1rem 0;
}

/* 파일 업로더 */
[data-testid="stFileUploader"] {
    background: #FFFFFF;
    border: 2px dashed #D1D5DB;
    border-radius: 12px;
    padding: 1rem;
}
[data-testid="stFileUploader"]:hover {
    border-color: #E61E3D;
}

/* 다운로드 버튼 */
.stDownloadButton button {
    background-color: #374151 !important;
    color: #FFFFFF !important;
    border: none !important;
    border-radius: 8px !important;
    font-weight: 600 !important;
}
.stDownloadButton button:hover {
    background-color: #4B5563 !important;
}

/* plotly 차트 배경 투명 */
.js-plotly-plot .plotly .main-svg {
    background: transparent !important;
}
</style>
""", unsafe_allow_html=True)

# ─── 메인 헤더 ───
st.markdown("""
<div class="main-header">
    <h1>COSMAX R&amp;I 매출분석</h1>
    <p>R&I Center · Sales Performance Analytics</p>
    <div class="accent-bar"></div>
</div>
""", unsafe_allow_html=True)

# ─── 파일 업로드 ───
uploaded_file = st.file_uploader("엑셀 파일을 업로드하세요", type=["xlsx", "xls"])
if uploaded_file is None:
    st.info("엑셀 파일을 업로드하면 분석이 시작됩니다.")
    st.stop()

# ─── 데이터 로드 ───
@st.cache_data
def load_data(file):
    xls = pd.ExcelFile(file)
    frames = []
    skip_keywords = ["요약", "작성", "주의", "매핑", "guide", "note"]
    for sheet in xls.sheet_names:
        if any(kw in sheet for kw in skip_keywords):
            continue
        df = pd.read_excel(xls, sheet_name=sheet)
        # 필수 컬럼 존재 여부 확인
        if "년도" not in df.columns or "월" not in df.columns:
            continue
        frames.append(df)
    data = pd.concat(frames, ignore_index=True)

    # 기본 전처리
    data["년도"] = pd.to_numeric(data["년도"], errors="coerce")
    data["월"] = pd.to_numeric(data["월"], errors="coerce")
    data = data.dropna(subset=["년도", "월"]).copy()
    data["년도"] = data["년도"].astype(int)
    data["월"] = data["월"].astype(int)
    data["년월"] = data["년도"].astype(str) + "-" + data["월"].astype(str).str.zfill(2)
    data["제품매출"] = pd.to_numeric(data["제품매출"], errors="coerce").fillna(0)
    data["순매출액"] = pd.to_numeric(data["순매출액"], errors="coerce").fillna(0)

    # 고객사 코드 추출 (상품코드 2~4번째 글자)
    data["고객코드"] = data["상품코드"].astype(str).str[1:4]

    # 고객사명 추출: 상품명에서 브랜드 부분 추출
    data["고객사명"] = data["상품코드"].astype(str).str[1:4]  # 기본값
    data = _map_customer_names(data)

    return data


def _map_customer_names(df):
    """상품명에서 고객사명을 추론하여 매핑"""
    customer_names = {}
    for code, group in df.groupby("고객코드"):
        names = group["상품명"].dropna().values
        if len(names) == 0:
            customer_names[code] = code
            continue
        # 상품명의 첫 번째 샘플에서 브랜드명 추출 시도
        sample = str(names[0])
        # 한글 + 영문 브랜드명 패턴 (보통 첫 단어가 브랜드)
        customer_names[code] = _extract_brand(sample, code)
    df["고객사명"] = df["고객코드"].map(customer_names)
    return df


def _extract_brand(product_name, code):
    """상품명에서 브랜드명 추출 - 알려진 브랜드 매핑 우선 적용"""
    known_brands = {
        "ABC": "미샤(에이블씨엔씨)",
        "DPD": "아누아",
        "SNT": "쏘내추럴",
        "OLV": "올리브영",
        "FCM": "어뮤즈/포컴퍼니",
        "BNU": "넘버즈인",
        "DPL": "닥터엘시아",
        "DKP": "센텔리안24(동국제약)",
        "AKC": "에스크컴퍼니",
        "SLC": "라운드랩",
        "CLO": "클리오/구달",
        "APS": "메디큐브(에이피알)",
        "GDI": "조선미녀",
        "CBK": "AHC",
        "DNC": "DNC(이지듀)",
        "JSM": "정샘물",
        "SKC": "스킨천사",
        "MDH": "메디힐",
        "GEN": "닥터지(고운세상)",
        "LOJ": "키엘(로레알)",
        "BOS": "이퀄(부스터스)",
        "PUR": "퍼셀",
        "NRP": "네이처리퍼블릭",
        "DNE": "닥터나인틴",
        "REB": "리파이",
        "TIR": "티르티르",
        "AHO": "앳홈",
        "HVB": "닥터자르트",
        "RJA": "닥터리쥬란",
        "ELC": "에스티로더",
        "TRD": "토리든",
        "ERK": "에런케이",
        "AWY": "암웨이",
        "MNF": "마녀공장",
        "CSB": "바이오던스",
        "DAP": "동아제약(파티온)",
        "AAL": "에이앤리더",
        "EZH": "닥터슈라이버(이지함)",
        "DTL": "닥터엘시아",
        "ENP": "홀리카홀리카",
        "CTO": "셀트리온스킨큐어",
        "IDP": "일동제약",
        "HJP": "휴젤(웰라쥬)",
        "FBL": "퓨쳐뷰티랩",
        "MLU": "멜라루카",
        "LAB": "레이어랩",
        "MIC": "메사니",
        "NOP": "네오팜",
        "FAK": "페이스리퍼블릭",
        "3AC": "신일제약",
        "MAM": "휴메딕스",
        "MNL": "에르보리앙",
        "KPE": "셀라듀",
        "BPT": "비플랜트",
        "BNC": "바닐라코",
        "HRM": "하루미",
        "KUL": "킬리킬리",
        "SSG": "연작(신세계)",
        "TNM": "토니모리",
        "MSA": "무신사",
        "SAM": "더샘",
        "EMT": "이마트",
        "WSN": "왓슨즈",
        "YHD": "유한킴벌리(그린핑거)",
        "YHK": "유한킴벌리",
        "HUX": "헉슬리",
        "KRT": "가히(코리아테크)",
        "MEM": "미미박스",
        "VIM": "바임",
        "PCT": "토코보(픽톤)",
        "LFC": "LF(아떼)",
        "FNF": "MLB(에프앤에프)",
        "SNG": "파뮤",
        "HMS": "한미사이언스",
    }
    if code in known_brands:
        return known_brands[code]
    return code


df = load_data(uploaded_file)


# ─── PPT 보고서 생성 ───
BLUE = RGBColor(30, 60, 120)
GRAY = RGBColor(100, 100, 100)
WHITE = RGBColor(255, 255, 255)
DARK = RGBColor(30, 30, 30)


def _save_chart(fig, dpi=180):
    """matplotlib figure -> 임시 파일 경로"""
    tmp = tempfile.NamedTemporaryFile(suffix=".png", delete=False)
    fig.savefig(tmp.name, dpi=dpi, bbox_inches="tight")
    plt.close(fig)
    return tmp.name


def _make_bar_chart(labels, values, title, ylabel="매출(억원)", color="#4472C4", figsize=(10, 4.5)):
    fig, ax = plt.subplots(figsize=figsize)
    bars = ax.bar(range(len(labels)), [v / 1e8 for v in values], color=color, width=0.6)
    ax.set_xticks(range(len(labels)))
    ax.set_xticklabels(labels, rotation=45, ha="right", fontsize=8)
    ax.set_ylabel(ylabel, fontsize=9)
    ax.set_title(title, fontsize=12, fontweight="bold", pad=12)
    for bar, v in zip(bars, values):
        ax.text(bar.get_x() + bar.get_width() / 2, bar.get_height(),
                f"{v/1e8:.1f}", ha="center", va="bottom", fontsize=7)
    ax.spines["top"].set_visible(False)
    ax.spines["right"].set_visible(False)
    plt.tight_layout()
    return _save_chart(fig)


def _make_line_chart(pivot_df, title, ylabel="매출(억원)", figsize=(10, 4.5)):
    fig, ax = plt.subplots(figsize=figsize)
    for col in pivot_df.columns:
        vals = pivot_df[col] / 1e8
        ax.plot(pivot_df.index, vals, marker="o", label=col, linewidth=1.5, markersize=5)
        for x, y in zip(pivot_df.index, vals):
            ax.text(x, y, f"{y:.1f}", fontsize=6, ha="center", va="bottom")
    ax.set_ylabel(ylabel, fontsize=9)
    ax.set_title(title, fontsize=12, fontweight="bold", pad=12)
    ax.legend(fontsize=7, loc="upper left", bbox_to_anchor=(1, 1))
    ax.spines["top"].set_visible(False)
    ax.spines["right"].set_visible(False)
    plt.tight_layout()
    return _save_chart(fig)


def _make_pie_chart(labels, values, title, figsize=(6, 4.5)):
    fig, ax = plt.subplots(figsize=figsize)
    ax.pie(values, labels=labels, autopct="%1.1f%%", startangle=90,
           textprops={"fontsize": 10})
    ax.set_title(title, fontsize=12, fontweight="bold", pad=12)
    plt.tight_layout()
    return _save_chart(fig)


def _add_slide_title(slide, title_text):
    """슬라이드 상단에 제목 텍스트박스 추가"""
    from pptx.util import Inches, Pt
    txBox = slide.shapes.add_textbox(Inches(0.5), Inches(0.2), Inches(9), Inches(0.6))
    tf = txBox.text_frame
    p = tf.paragraphs[0]
    p.text = title_text
    p.font.size = Pt(22)
    p.font.bold = True
    p.font.color.rgb = BLUE


def _add_table_to_slide(slide, headers, rows, left, top, width, height):
    """슬라이드에 테이블 추가"""
    from pptx.util import Inches, Pt, Emu
    n_rows = len(rows) + 1
    n_cols = len(headers)
    table_shape = slide.shapes.add_table(n_rows, n_cols, left, top, width, height)
    table = table_shape.table

    col_w = int(width / n_cols)
    for i in range(n_cols):
        table.columns[i].width = col_w

    # 헤더
    for i, h in enumerate(headers):
        cell = table.cell(0, i)
        cell.text = str(h)
        for paragraph in cell.text_frame.paragraphs:
            paragraph.font.size = Pt(9)
            paragraph.font.bold = True
            paragraph.font.color.rgb = WHITE
            paragraph.alignment = PP_ALIGN.CENTER
        cell.fill.solid()
        cell.fill.fore_color.rgb = BLUE

    # 데이터
    for ri, row in enumerate(rows):
        for ci, val in enumerate(row):
            cell = table.cell(ri + 1, ci)
            cell.text = str(val)
            for paragraph in cell.text_frame.paragraphs:
                paragraph.font.size = Pt(8)
                paragraph.font.color.rgb = DARK
                paragraph.alignment = PP_ALIGN.CENTER
            if ri % 2 == 0:
                cell.fill.solid()
                cell.fill.fore_color.rgb = RGBColor(230, 236, 245)


def generate_pptx(data, metric):
    """PPT 보고서 생성 (네이티브 차트 - PPT에서 편집 가능) -> bytes 반환"""
    from pptx.chart.data import CategoryChartData
    template_path = os.path.join(os.path.dirname(__file__), "2026년 사업계획 양식.pptx")

    # 템플릿 로드
    if os.path.exists(template_path):
        prs = Presentation(template_path)
        while len(prs.slides) > 0:
            rId = prs.slides._sldIdLst[0].rId
            prs.part.drop_rel(rId)
            del prs.slides._sldIdLst[0]
    else:
        prs = Presentation()
        prs.slide_width = Emu(12192000)
        prs.slide_height = Emu(6858000)

    content_layout = prs.slide_layouts[6]  # 빈 화면
    cover_layout = end_layout = None
    for layout in prs.slide_layouts:
        if "표지" in layout.name or "Title" in layout.name:
            cover_layout = layout
        elif "End" in layout.name or "end" in layout.name:
            end_layout = layout
    if cover_layout is None:
        cover_layout = content_layout

    months = sorted(data["년월"].unique())
    period = f"{months[0]} ~ {months[-1]}" if len(months) > 1 else months[0]
    today = datetime.date.today().strftime("%Y.%m.%d")

    def _new_slide(title_text):
        s = prs.slides.add_slide(content_layout)
        s.background.fill.solid()
        s.background.fill.fore_color.rgb = RGBColor(255, 255, 255)
        ln = s.shapes.add_connector(1, Emu(347134), Emu(745066), Emu(11709400), Emu(745066))
        ln.line.color.rgb = RGBColor(200, 200, 200)
        ln.line.width = Pt(0.75)
        tx = s.shapes.add_textbox(Emu(399560), Emu(254430), Emu(8000000), Emu(400110))
        p = tx.text_frame.paragraphs[0]
        p.text = title_text
        p.font.size = Pt(20)
        p.font.bold = True
        p.font.color.rgb = RGBColor(0, 0, 0)
        tag = s.shapes.add_shape(1, Emu(10315094), Emu(397263), Emu(1409700), Emu(323850))
        tag.fill.solid()
        tag.fill.fore_color.rgb = RGBColor(0, 112, 192)
        tag.line.fill.background()
        tp = tag.text_frame.paragraphs[0]
        tp.text = "R&I"
        tp.font.size = Pt(12)
        tp.font.bold = True
        tp.font.color.rgb = WHITE
        tp.alignment = PP_ALIGN.CENTER
        return s

    def _add_bar(slide, labels, values, title, left, top, width, height, color_hex="0070C0"):
        cd = CategoryChartData()
        cd.categories = labels
        cd.add_series("매출(억원)", [v / 1e8 for v in values])
        cf = slide.shapes.add_chart(XL_CHART_TYPE.COLUMN_CLUSTERED, left, top, width, height, cd)
        ch = cf.chart
        ch.has_legend = False
        ch.chart_title.has_text_frame = True
        ch.chart_title.text_frame.paragraphs[0].text = title
        ch.chart_title.text_frame.paragraphs[0].font.size = Pt(11)
        ch.chart_title.text_frame.paragraphs[0].font.bold = True
        pl = ch.plots[0]
        pl.gap_width = 80
        sr = pl.series[0]
        sr.format.fill.solid()
        sr.format.fill.fore_color.rgb = RGBColor.from_string(color_hex)
        sr.data_labels.show_value = True
        sr.data_labels.font.size = Pt(7)
        sr.data_labels.number_format = '#,##0.0'

    def _add_line(slide, pivot_df, title, left, top, width, height):
        cd = CategoryChartData()
        cd.categories = pivot_df.index.tolist()
        for col in pivot_df.columns:
            cd.add_series(str(col), [v / 1e8 for v in pivot_df[col].values])
        cf = slide.shapes.add_chart(XL_CHART_TYPE.LINE_MARKERS, left, top, width, height, cd)
        ch = cf.chart
        ch.has_legend = True
        ch.legend.position = XL_LEGEND_POSITION.BOTTOM
        ch.legend.font.size = Pt(7)
        ch.chart_title.has_text_frame = True
        ch.chart_title.text_frame.paragraphs[0].text = title
        ch.chart_title.text_frame.paragraphs[0].font.size = Pt(11)
        for s in ch.plots[0].series:
            s.data_labels.show_value = True
            s.data_labels.font.size = Pt(6)
            s.data_labels.number_format = '#,##0.0'

    def _add_pie(slide, labels, values, title, left, top, width, height):
        cd = CategoryChartData()
        cd.categories = labels
        cd.add_series("매출", [v / 1e8 for v in values])
        cf = slide.shapes.add_chart(XL_CHART_TYPE.PIE, left, top, width, height, cd)
        ch = cf.chart
        ch.has_legend = True
        ch.legend.position = XL_LEGEND_POSITION.BOTTOM
        ch.legend.font.size = Pt(8)
        ch.chart_title.has_text_frame = True
        ch.chart_title.text_frame.paragraphs[0].text = title
        ch.chart_title.text_frame.paragraphs[0].font.size = Pt(11)
        ch.plots[0].series[0].data_labels.show_percentage = True
        ch.plots[0].series[0].data_labels.font.size = Pt(8)

    def _add_grouped_bar(slide, pivot_df, title, left, top, width, height, colors=None):
        cd = CategoryChartData()
        cd.categories = pivot_df.index.tolist()
        for col in pivot_df.columns:
            cd.add_series(str(col), [v / 1e8 for v in pivot_df[col].values])
        cf = slide.shapes.add_chart(XL_CHART_TYPE.COLUMN_CLUSTERED, left, top, width, height, cd)
        ch = cf.chart
        ch.has_legend = True
        ch.legend.position = XL_LEGEND_POSITION.BOTTOM
        ch.chart_title.has_text_frame = True
        ch.chart_title.text_frame.paragraphs[0].text = title
        ch.chart_title.text_frame.paragraphs[0].font.size = Pt(11)
        if colors:
            for i, s in enumerate(ch.plots[0].series):
                if i < len(colors):
                    s.format.fill.solid()
                    s.format.fill.fore_color.rgb = RGBColor.from_string(colors[i])

    # 슬라이드 1: 표지
    slide = prs.slides.add_slide(cover_layout)
    for shape in slide.shapes:
        if shape.is_placeholder:
            idx = shape.placeholder_format.idx
            if idx == 11:
                for p in shape.text_frame.paragraphs:
                    for r in p.runs:
                        if "20XX" in r.text: r.text = today
                        elif "Lab" in r.text: r.text = "ES Lab"
            elif idx == 12:
                for p in shape.text_frame.paragraphs:
                    for r in p.runs:
                        if "COSMAX" in r.text: r.text = "COSMAX R&I"
                        elif "매출" in r.text: r.text = f"매출 분석 자료 ({period})"

    # 슬라이드 2: 목차
    slide = _new_slide("목 차")
    for idx, item in enumerate(["I. 실적 현황 요약", "II. 고객사별 매출 분석", "III. 품목(유형)별 매출 분석",
                                 "IV. 국내/해외 매출 분석", "V. 담당자별 매출 분석", "VI. 키워드 트렌드 분석"]):
        tx = slide.shapes.add_textbox(Emu(1468487), Emu(1100000 + idx * 650000), Emu(6000000), Emu(400110))
        p = tx.text_frame.paragraphs[0]
        p.text = item
        p.font.size = Pt(18)
        p.font.bold = True

    # 슬라이드 3: 실적 현황 요약
    slide = _new_slide("I. 실적 현황 요약")
    total = data[metric].sum()
    dom_total = data[data["국내/해외"] == "국내"][metric].sum()
    ovs_total = data[data["국내/해외"] == "해외"][metric].sum()
    _add_table_to_slide(slide, ["구분", "금액(억원)", "비중"], [
        ["총 매출", f"{total/1e8:,.1f}", "100.0%"],
        ["국내 매출", f"{dom_total/1e8:,.1f}", f"{dom_total/total*100:.1f}%" if total else "0%"],
        ["해외 매출", f"{ovs_total/1e8:,.1f}", f"{ovs_total/total*100:.1f}%" if total else "0%"],
        ["고객사 수", f"{data['고객코드'].nunique()}개", "-"],
        ["상품 수", f"{data['상품코드'].nunique()}개", "-"],
    ], Inches(0.5), Inches(1.0), Inches(5), Inches(2.2))
    team_sum = data.groupby("Team")[metric].sum().sort_values(ascending=False).reset_index()
    _add_table_to_slide(slide, ["Team", "매출(억원)", "비중"],
        [[r["Team"], f"{r[metric]/1e8:,.1f}", f"{r[metric]/total*100:.1f}%" if total else "0%"] for _, r in team_sum.iterrows()],
        Inches(6.5), Inches(1.0), Inches(5), Inches(1.5))

    # 슬라이드 3-1: 실적 현황표
    slide = _new_slide("I-1. 고객사별 실적 현황표")
    _ms = sorted(data["년월"].unique())
    _pv = data.pivot_table(index="고객사명", columns="년월", values=metric, aggfunc="sum").fillna(0)
    _pv = _pv.reindex(columns=_ms, fill_value=0)
    _pv["누적"] = _pv.sum(axis=1)
    _pv["제품수"] = data.groupby("고객사명")["상품코드"].nunique()
    _pv = _pv.sort_values("누적", ascending=False).head(20)
    _gt = _pv["누적"].sum()
    th = ["분류"]
    for i, m in enumerate(_ms):
        th.append(metric)
        if i > 0: th.extend(["증감", "성장률"])
    th.extend(["누적", "비중", "순위", "제품수"])
    tr0 = ["총합계"]
    for i, m in enumerate(_ms):
        v = data[data["년월"] == m][metric].sum()
        tr0.append(f"{v/1e8:,.1f}억")
        if i > 0:
            pv = data[data["년월"] == _ms[i-1]][metric].sum()
            d = v - pv; rt = (d/pv*100) if pv else 0; a = "▲" if d >= 0 else "▼"
            tr0.append(f"{a} {abs(d)/1e8:,.1f}억"); tr0.append(f"{a} {abs(rt):.1f}%")
    tr0.extend([f"{_gt/1e8:,.1f}억", "100%", "-", str(data["상품코드"].nunique())])
    trs = [tr0]
    for rk, (nm, rw) in enumerate(_pv.iterrows(), 1):
        r = [nm]
        for i, m in enumerate(_ms):
            v = rw[m]; r.append(f"{v/1e8:,.1f}억")
            if i > 0:
                pv = rw[_ms[i-1]]; d = v-pv; rt = (d/pv*100) if pv else 0; a = "▲" if d >= 0 else "▼"
                r.append(f"{a} {abs(d)/1e8:,.1f}억"); r.append(f"{a} {abs(rt):.1f}%")
        cm = rw["누적"]
        r.extend([f"{cm/1e8:,.1f}억", f"{cm/_gt*100:.1f}%" if _gt else "0%", str(rk), str(int(rw["제품수"]))])
        trs.append(r)
    _add_table_to_slide(slide, th, trs, Emu(200000), Emu(1000000), Emu(11700000), Emu(min(len(trs)*200000+300000, 5500000)))

    # 슬라이드 4: 고객사별 매출 순위
    slide = _new_slide("II. 고객사별 매출 순위 (Top 20)")
    cust = data.groupby("고객사명")[metric].sum().sort_values(ascending=False).head(20)
    _add_bar(slide, cust.index.tolist(), cust.values.tolist(), "고객사별 매출 Top 20 (억원)",
             Emu(300000), Emu(1100000), Emu(7800000), Emu(3800000))
    cust_tbl = data.groupby(["고객사명", "고객코드"]).agg(
        매출합계=(metric, "sum"), 상품수=("상품코드", "nunique")
    ).sort_values("매출합계", ascending=False).head(10).reset_index()
    _add_table_to_slide(slide, ["#", "고객사", "매출(억)", "상품수"],
        [[i+1, r["고객사명"][:8], f"{r['매출합계']/1e8:,.1f}", r["상품수"]] for i, (_, r) in enumerate(cust_tbl.iterrows())],
        Emu(8300000), Emu(1100000), Emu(3600000), Emu(3800000))

    # 슬라이드 5: 고객사별 월별 추이
    slide = _new_slide("II-1. 고객사별 월별 매출 추이 (Top 10)")
    top10 = data.groupby("고객사명")[metric].sum().sort_values(ascending=False).head(10).index
    pv = data[data["고객사명"].isin(top10)].pivot_table(index="년월", columns="고객사명", values=metric, aggfunc="sum").fillna(0)
    pv = pv[top10]
    _add_line(slide, pv, "고객사별 월별 매출 추이 (억원)", Emu(300000), Emu(1100000), Emu(11500000), Emu(5200000))

    # 슬라이드 6: 품목별
    slide = _new_slide("III. 품목(유형)별 매출 분석")
    ts = data.groupby("중유형")[metric].sum().sort_values(ascending=False)
    _add_bar(slide, ts.index.tolist(), ts.values.tolist(), "중유형별 매출 (억원)",
             Emu(300000), Emu(1100000), Emu(5700000), Emu(2900000))
    ss = data.groupby("소유형")[metric].sum().sort_values(ascending=False).head(12)
    _add_bar(slide, ss.index.tolist(), ss.values.tolist(), "소유형별 매출 Top 12 (억원)",
             Emu(6200000), Emu(1100000), Emu(5700000), Emu(2900000), "374151")
    cs = data.groupby("19년 카테고리")[metric].sum().sort_values(ascending=False)
    _add_pie(slide, cs.index.tolist(), cs.values.tolist(), "카테고리별 매출 비중",
             Emu(3200000), Emu(4100000), Emu(5500000), Emu(2600000))

    # 슬라이드 7: 국내/해외
    slide = _new_slide("IV. 국내/해외 매출 분석")
    rg = data.groupby("국내/해외")[metric].sum()
    _add_pie(slide, rg.index.tolist(), rg.values.tolist(), "국내 vs 해외 매출 비중",
             Emu(300000), Emu(1100000), Emu(4200000), Emu(3200000))
    rt = data.groupby(["국내/해외", "중유형"])[metric].sum().reset_index()
    pvt = rt.pivot_table(index="중유형", columns="국내/해외", values=metric, aggfunc="sum").fillna(0)
    _add_grouped_bar(slide, pvt, "국내/해외별 중유형 비교 (억원)",
                     Emu(4800000), Emu(1100000), Emu(5100000), Emu(3200000), colors=["0070C0", "E61E3D"])

    # 슬라이드 8: 국내/해외 Top10
    slide = _new_slide("IV-1. 국내/해외 고객사 매출 Top 10")
    for idx, (rn, clr) in enumerate([("국내", "0070C0"), ("해외", "E61E3D")]):
        rd = data[data["국내/해외"] == rn]
        if len(rd) == 0: continue
        rc = rd.groupby("고객사명")[metric].sum().sort_values(ascending=False).head(10)
        xp = Emu(300000) if idx == 0 else Emu(6200000)
        _add_bar(slide, rc.index.tolist(), rc.values.tolist(), f"{rn} 고객사 Top 10 (억원)",
                 xp, Emu(1100000), Emu(5700000), Emu(5000000), clr)

    # 슬라이드 9: 담당자별
    slide = _new_slide("V. 담당자별 매출 실적")
    mg = data.groupby(["현담당자", "현담당자 팀"]).agg(
        매출합계=(metric, "sum"), 담당제품수=("상품코드", "nunique"), 담당고객수=("고객코드", "nunique")
    ).sort_values("매출합계", ascending=False).reset_index()
    _add_bar(slide, mg["현담당자"].tolist(), mg["매출합계"].tolist(), "담당자별 매출 순위 (억원)",
             Emu(300000), Emu(1100000), Emu(7000000), Emu(3200000))
    _add_table_to_slide(slide, ["#", "담당자", "팀", "매출(억)", "제품수", "고객수"],
        [[i+1, r["현담당자"], r["현담당자 팀"], f"{r['매출합계']/1e8:,.1f}", r["담당제품수"], r["담당고객수"]]
         for i, (_, r) in enumerate(mg.iterrows())],
        Emu(7500000), Emu(1100000), Emu(4400000), Emu(min(len(mg)*180000+300000, 5000000)))

    # 슬라이드 10: 키워드
    slide = _new_slide("VI. 제품명 키워드 트렌드 분석")
    _KW = ["미백","브라이트","잡티","주름","탄력","리프트","보습","수분","히알루론",
           "진정","시카","카밍","각질","필링","모공","비타민","비타","레티놀",
           "콜라겐","펩타이드","세라마이드","프로폴리스","병풀","마데카",
           "앰플","세럼","에센스","토너","크림","마스크","패드","미스트"]
    kw_rec = []
    for _, row in data.iterrows():
        nm = str(row["상품명"]).upper()
        for kw in _KW:
            if kw.upper() in nm:
                kw_rec.append({"키워드": kw, metric: row[metric]})
    if kw_rec:
        kwa = pd.DataFrame(kw_rec).groupby("키워드")[metric].sum().sort_values(ascending=False).head(10)
        _add_bar(slide, kwa.index.tolist(), kwa.values.tolist(), "제품명 키워드 매출 Top 10 (억원)",
                 Emu(300000), Emu(1100000), Emu(11500000), Emu(5200000))

    # 마지막: 엔딩
    if end_layout:
        prs.slides.add_slide(end_layout)

    buf = io.BytesIO()
    prs.save(buf)
    buf.seek(0)
    return buf.getvalue()

# ─── 사이드바 필터 ───
st.sidebar.header("🔍 필터")

# 국내/해외 필터
region_options = ["전체"] + sorted(df["국내/해외"].unique().tolist())
selected_region = st.sidebar.selectbox("국내/해외", region_options)

# 년도 필터
year_options = ["전체"] + sorted(df["년도"].unique().astype(str).tolist())
selected_year = st.sidebar.selectbox("년도", year_options)

# 분기 필터
df["분기"] = "Q" + ((df["월"] - 1) // 3 + 1).astype(str)
quarter_options = ["전체"] + sorted(df["분기"].unique().tolist())
selected_quarter = st.sidebar.selectbox("분기", quarter_options)

# 월 필터 (멀티셀렉트)
all_months = sorted(df["년월"].unique().tolist())
selected_months = st.sidebar.multiselect("월 (복수 선택 가능)", all_months, default=all_months)

# Team 필터
team_options = ["전체"] + sorted(df["Team"].dropna().unique().tolist())
selected_team = st.sidebar.selectbox("Team", team_options)

# 매출 기준 선택
metric = st.sidebar.radio("매출 기준", ["제품매출", "순매출액"])

# 상위 고객사 수
top_n = st.sidebar.slider("상위 고객사 수", min_value=5, max_value=50, value=10)

# 필터 적용
filtered = df.copy()
if selected_region != "전체":
    filtered = filtered[filtered["국내/해외"] == selected_region]
if selected_year != "전체":
    filtered = filtered[filtered["년도"] == int(selected_year)]
if selected_quarter != "전체":
    filtered = filtered[filtered["분기"] == selected_quarter]
if selected_months:
    filtered = filtered[filtered["년월"].isin(selected_months)]
else:
    filtered = filtered[filtered["년월"].isin(all_months)]
if selected_team != "전체":
    filtered = filtered[filtered["Team"] == selected_team]

# ─── PPT 보고서 다운로드 ───
st.sidebar.divider()
st.sidebar.header("📄 보고서 다운로드")
if st.sidebar.button("PPT 보고서 생성", type="primary", use_container_width=True):
    with st.sidebar.status("PPT 생성 중...", expanded=True) as status:
        pptx_bytes = generate_pptx(filtered, metric)
        status.update(label="PPT 생성 완료!", state="complete")
    months = sorted(filtered["년월"].unique())
    period_start = months[0].replace("-", "") if months else ""
    period_end = months[-1].replace("-", "") if months else ""
    today = datetime.date.today().strftime("%Y%m%d")
    filename = f"현황보고 매출 실적 ({period_start}-{period_end}, {today}).pptx"
    st.sidebar.download_button(
        "📥 PPT 다운로드",
        data=pptx_bytes,
        file_name=filename,
        mime="application/vnd.openxmlformats-officedocument.presentationml.presentation",
        use_container_width=True,
    )

# ─── 사이드바 하단 크레딧 ───
st.sidebar.divider()
st.sidebar.markdown(
    f"""<div style="text-align:center; padding:8px 0 12px 0;">
    <span style="color:#6B7280; font-size:10px; letter-spacing:1.5px; text-transform:uppercase;">
    Uploaded {datetime.date.today().strftime('%b %d, %Y')}</span><br>
    <span style="color:#9CA3AF; font-size:9px; letter-spacing:2px;">
    Designed by <b style="color:#E5E7EB;">CHOI SUJUNG</b></span>
    </div>""",
    unsafe_allow_html=True,
)

# ─── 상단 KPI ───
_filtered_months = sorted(filtered["년월"].unique())
_period_str = f"{_filtered_months[0]} ~ {_filtered_months[-1]}" if len(_filtered_months) > 1 else (_filtered_months[0] if _filtered_months else "-")
st.caption(f"분석 기간: **{_period_str}** ({len(_filtered_months)}개월) | 데이터: **{len(filtered):,}**건")

col1, col2, col3, col4, col5 = st.columns(5)
total_sales = filtered[metric].sum()
domestic = filtered[filtered["국내/해외"] == "국내"][metric].sum()
overseas = filtered[filtered["국내/해외"] == "해외"][metric].sum()
n_customers = filtered["고객코드"].nunique()
n_products = filtered["상품코드"].nunique()

col1.metric("총 매출", f"{total_sales/1e8:,.1f}억원")
col2.metric("국내 매출", f"{domestic/1e8:,.1f}억원")
col3.metric("해외 매출", f"{overseas/1e8:,.1f}억원")
col4.metric("고객사 수", f"{n_customers}개")
col5.metric("상품 수", f"{n_products}개")

st.divider()

# ─── Plotly 공통 스타일 ───
PI_COLORS = ["#E61E3D", "#1F2937", "#6B7280", "#9CA3AF", "#374151",
             "#DC2626", "#4B5563", "#F87171", "#111827", "#D1D5DB"]
PI_GRADIENT = [[0, "#374151"], [0.5, "#6B7280"], [1, "#E61E3D"]]

def _pi_layout(fig, **kwargs):
    """PI 스타일 공통 레이아웃 적용"""
    fig.update_layout(
        plot_bgcolor="rgba(0,0,0,0)",
        paper_bgcolor="rgba(0,0,0,0)",
        font=dict(family="Noto Sans KR, Malgun Gothic", color="#1F2937"),
        **kwargs,
    )
    return fig

def _excel_download(dataframe, filename, key):
    """DataFrame을 엑셀로 다운로드하는 버튼 생성"""
    buf = io.BytesIO()
    with pd.ExcelWriter(buf, engine="openpyxl") as writer:
        dataframe.to_excel(writer, index=False, sheet_name="data")
    st.download_button(
        "📥 엑셀 다운로드",
        data=buf.getvalue(),
        file_name=f"{filename}.xlsx",
        mime="application/vnd.openxmlformats-officedocument.spreadsheetml.sheet",
        key=key,
    )

# ─── 탭 구성 ───
tab0, tab1, tab2, tab3, tab5, tab7, tab6, tab4, tab_help = st.tabs(
    ["📊 실적 현황표", "🏢 고객사 분석", "📦 품목(유형) 분석", "🌏 국내/해외 분석", "👤 담당자 분석", "🔑 키워드 분석", "🔍 인터랙티브 분석", "📋 상세 데이터", "📖 작성 방법"]
)

# ============================================================
# 탭0: 실적 현황표
# ============================================================
with tab0:

    def _arrow(val):
        """증감에 따른 화살표 + 색상 HTML"""
        if val > 0:
            return f'<span style="color:#E61E3D;font-weight:700">▲ {abs(val)/1e8:,.1f}억</span>'
        elif val < 0:
            return f'<span style="color:#2563EB;font-weight:700">▼ {abs(val)/1e8:,.1f}억</span>'
        return '<span style="color:#6B7280">- 0</span>'

    def _rate_arrow(val):
        if val > 0:
            return f'<span style="color:#E61E3D;font-weight:700">▲ {val:.1f}%</span>'
        elif val < 0:
            return f'<span style="color:#2563EB;font-weight:700">▼ {abs(val):.1f}%</span>'
        return '<span style="color:#6B7280">0%</span>'

    def _amt(val):
        return f'{val/1e8:,.1f}억'

    months_sorted = sorted(df["년월"].unique())
    report_view = st.radio(
        "분석 기준 선택",
        ["고객사별 실적", "중유형별 실적", "담당자별 실적", "Team별 실적"],
        horizontal=True,
        key="report_view",
    )

    view_col_map = {
        "고객사별 실적": "고객사명",
        "중유형별 실적": "중유형",
        "담당자별 실적": "현담당자",
        "Team별 실적": "Team",
    }
    group_col = view_col_map[report_view]

    # 월별 데이터 피벗
    monthly_pivot = df.pivot_table(
        index=group_col, columns="년월", values=metric, aggfunc="sum"
    ).fillna(0)
    monthly_pivot = monthly_pivot.reindex(columns=months_sorted, fill_value=0)

    # 누적 합계
    monthly_pivot["누적합계"] = monthly_pivot.sum(axis=1)

    # 제품수
    prod_count = df.groupby(group_col)["상품코드"].nunique()
    monthly_pivot["제품수"] = prod_count

    # 정렬
    monthly_pivot = monthly_pivot.sort_values("누적합계", ascending=False)

    # 전체 합계 행
    total_row = monthly_pivot.sum()
    total_row.name = "총합계"

    # Top N 적용
    display_pivot = monthly_pivot.head(top_n)

    # HTML 테이블 생성
    html = """
    <style>
    .report-table {
        width: 100%;
        border-collapse: collapse;
        font-family: 'Noto Sans KR', 'Malgun Gothic', sans-serif;
        font-size: 13px;
    }
    .report-table th {
        background: #111827;
        color: #FFFFFF;
        padding: 8px 12px;
        text-align: center;
        border: 1px solid #374151;
        font-weight: 700;
        font-size: 12px;
        white-space: nowrap;
    }
    .report-table td {
        padding: 6px 10px;
        border: 1px solid #E5E7EB;
        text-align: right;
        white-space: nowrap;
    }
    .report-table tr:nth-child(even) {
        background: #F9FAFB;
    }
    .report-table tr:hover {
        background: #F0F4FA;
    }
    .report-table .row-name {
        text-align: left;
        font-weight: 600;
        color: #1F2937;
        background: #F3F4F6;
        position: sticky;
        left: 0;
    }
    .report-table .total-row td {
        background: #1F2937 !important;
        color: #FFFFFF !important;
        font-weight: 700;
    }
    .report-table .total-row span {
        color: #FCA5A5 !important;
    }
    </style>
    <div style="overflow-x:auto;">
    <table class="report-table">
    <thead>
    <tr>
        <th rowspan="2" style="min-width:120px">분류</th>
    """

    # 월별 헤더 (각 월마다 실적 + 전월대비)
    for i, m in enumerate(months_sorted):
        if i == 0:
            html += f'<th colspan="1">{m}</th>'
        else:
            html += f'<th colspan="3">{m}</th>'

    html += '<th colspan="3">누적</th>'
    html += '<th rowspan="2">제품수</th>'
    html += '</tr><tr>'

    for i, m in enumerate(months_sorted):
        if i == 0:
            html += f'<th>{metric}</th>'
        else:
            html += f'<th>{metric}</th><th>증감</th><th>성장률</th>'

    html += f'<th>{metric}</th><th>비중</th><th>순위</th>'
    html += '</tr></thead><tbody>'

    grand_total = total_row["누적합계"]

    # 총합계 행
    html += '<tr class="total-row">'
    html += f'<td class="row-name" style="background:#1F2937!important;color:#FFF!important">총합계</td>'
    for i, m in enumerate(months_sorted):
        val = total_row[m]
        if i == 0:
            html += f'<td>{_amt(val)}</td>'
        else:
            prev_m = months_sorted[i - 1]
            prev_val = total_row[prev_m]
            diff = val - prev_val
            rate = (diff / prev_val * 100) if prev_val else 0
            html += f'<td>{_amt(val)}</td><td>{_arrow(diff)}</td><td>{_rate_arrow(rate)}</td>'
    html += f'<td>{_amt(grand_total)}</td><td>100%</td><td>-</td>'
    html += f'<td>{int(total_row["제품수"])}</td>'
    html += '</tr>'

    # 각 행
    for rank, (name, row) in enumerate(display_pivot.iterrows(), 1):
        html += '<tr>'
        html += f'<td class="row-name">{name}</td>'
        for i, m in enumerate(months_sorted):
            val = row[m]
            if i == 0:
                html += f'<td>{_amt(val)}</td>'
            else:
                prev_m = months_sorted[i - 1]
                prev_val = row[prev_m]
                diff = val - prev_val
                rate = (diff / prev_val * 100) if prev_val else 0
                html += f'<td>{_amt(val)}</td><td>{_arrow(diff)}</td><td>{_rate_arrow(rate)}</td>'

        cum = row["누적합계"]
        share = (cum / grand_total * 100) if grand_total else 0
        html += f'<td style="font-weight:700">{_amt(cum)}</td>'
        html += f'<td>{share:.1f}%</td>'
        html += f'<td style="text-align:center;font-weight:700">{rank}</td>'
        html += f'<td style="text-align:center">{int(row["제품수"])}</td>'
        html += '</tr>'

    html += '</tbody></table></div>'

    st.subheader("실적 현황표")
    st.markdown(html, unsafe_allow_html=True)

    # 다운로드
    export_df = display_pivot.reset_index()
    col_dl1, col_dl2 = st.columns(2)
    with col_dl1:
        _excel_download(export_df, "실적현황표", "dl_report_table")
    with col_dl2:
        def _gen_report_ppt():
            _prs = Presentation()
            _prs.slide_width = Emu(12192000)
            _prs.slide_height = Emu(6858000)
            _sl = _prs.slides.add_slide(_prs.slide_layouts[6])
            _sl.background.fill.solid()
            _sl.background.fill.fore_color.rgb = RGBColor(255,255,255)
            _tx = _sl.shapes.add_textbox(Emu(300000), Emu(200000), Emu(8000000), Emu(450000))
            _p = _tx.text_frame.paragraphs[0]
            _p.text = f"실적 현황표 ({report_view})"
            _p.font.size = Pt(24)
            _p.font.bold = True
            _p.font.color.rgb = RGBColor(17,24,39)
            _ms = sorted(filtered["년월"].unique())
            _pv = filtered.pivot_table(index=group_col, columns="년월", values=metric, aggfunc="sum").fillna(0)
            _pv = _pv.reindex(columns=_ms, fill_value=0)
            _pv["누적"] = _pv.sum(axis=1)
            _pv["제품수"] = filtered.groupby(group_col)["상품코드"].nunique()
            _pv = _pv.sort_values("누적", ascending=False).head(top_n)
            _gt = _pv["누적"].sum()
            nmc = len(_ms) + max(0, (len(_ms)-1)*2)
            nc = 1 + nmc + 4
            nr = 2 + 1 + len(_pv)
            ts = _sl.shapes.add_table(nr, nc, Emu(150000), Emu(750000), Emu(11900000), Emu(min(nr*280000,5800000)))
            t = ts.table
            H=RGBColor(17,24,39); T=RGBColor(31,41,55); R=RGBColor(230,30,61); B=RGBColor(37,99,235)
            D=RGBColor(17,24,39); W=RGBColor(255,255,255); S=RGBColor(243,244,246)
            def _c(r,c,tx,b=False,bg=None,fg=None,sz=8):
                cl=t.cell(r,c); cl.text=str(tx)
                for p in cl.text_frame.paragraphs:
                    p.font.size=Pt(sz); p.font.bold=b; p.alignment=PP_ALIGN.CENTER
                    if fg: p.font.color.rgb=fg
                if bg: cl.fill.solid(); cl.fill.fore_color.rgb=bg
            _c(0,0,"분류",True,H,W,9)
            try: t.cell(0,0).merge(t.cell(1,0))
            except: pass
            ci=1
            for i,m in enumerate(_ms):
                if i==0: _c(0,ci,m,True,H,W,9); ci+=1
                else:
                    try: t.cell(0,ci).merge(t.cell(0,ci+2))
                    except: pass
                    _c(0,ci,m,True,H,W,9); ci+=3
            try: t.cell(0,ci).merge(t.cell(0,ci+2))
            except: pass
            _c(0,ci,"누적",True,H,W,9)
            _c(0,nc-1,"제품수",True,H,W,9)
            try: t.cell(0,nc-1).merge(t.cell(1,nc-1))
            except: pass
            ci=1
            for i,m in enumerate(_ms):
                _c(1,ci,metric,True,H,W,7); ci+=1
                if i>0: _c(1,ci,"증감",True,H,W,7); _c(1,ci+1,"성장률",True,H,W,7); ci+=2
            _c(1,ci,metric,True,H,W,7); _c(1,ci+1,"비중",True,H,W,7); _c(1,ci+2,"순위",True,H,W,7)
            def _r(ri,nm,mv,cum,sh,rk,pc,tot=False):
                bg=T if tot else (S if ri%2==0 else None); fg=W if tot else D
                _c(ri,0,nm,True,bg,fg,9); ci=1
                for i,m in enumerate(_ms):
                    v=mv.get(m,0); _c(ri,ci,f"{v/1e8:,.1f}억",tot,bg,fg,8); ci+=1
                    if i>0:
                        pv=mv.get(_ms[i-1],0); d=v-pv; rt=(d/pv*100) if pv else 0; a="▲" if d>=0 else "▼"
                        cl=(RGBColor(252,165,165) if d>=0 else RGBColor(147,197,253)) if tot else (R if d>=0 else B)
                        _c(ri,ci,f"{a} {abs(d)/1e8:,.1f}억",True,bg,cl,8); _c(ri,ci+1,f"{a} {abs(rt):.1f}%",True,bg,cl,8); ci+=2
                _c(ri,ci,f"{cum/1e8:,.1f}억",True,bg,fg,8); _c(ri,ci+1,f"{sh:.1f}%",False,bg,fg,8)
                _c(ri,ci+2,str(rk),True,bg,fg,8); _c(ri,nc-1,str(pc),False,bg,fg,8)
            tv={m:filtered[filtered["년월"]==m][metric].sum() for m in _ms}
            _r(2,"총합계",tv,_gt,100.0,"-",filtered["상품코드"].nunique(),True)
            for rk,(nm,rw) in enumerate(_pv.iterrows(),1):
                _r(3+rk-1,nm,{m:rw[m] for m in _ms},rw["누적"],(rw["누적"]/_gt*100) if _gt else 0,rk,int(rw["제품수"]))
            buf=io.BytesIO(); _prs.save(buf); buf.seek(0); return buf.getvalue()

        if st.button("📥 실적 현황표 PPT", key="dl_rpt_pptx"):
            ppt_data = _gen_report_ppt()
            st.download_button("💾 PPT 저장", data=ppt_data,
                file_name=f"실적현황표_{report_view}.pptx",
                mime="application/vnd.openxmlformats-officedocument.presentationml.presentation",
                key="dl_rpt_pptx_f")

# ============================================================
# 탭1: 고객사 분석
# ============================================================
with tab1:
    st.subheader(f"고객사별 매출 순위 (Top {top_n})")

    cust_sales = (
        filtered.groupby("고객사명")[metric]
        .sum()
        .sort_values(ascending=False)
        .head(top_n)
        .reset_index()
    )
    fig1 = px.bar(
        cust_sales,
        x="고객사명",
        y=metric,
        text_auto=".3s",
        color=metric,
        color_continuous_scale=[[0, "#374151"], [0.5, "#6B7280"], [1, "#E61E3D"]],
    )
    fig1.update_layout(
        xaxis_tickangle=-45, height=500, coloraxis_showscale=False,
        plot_bgcolor="rgba(0,0,0,0)", paper_bgcolor="rgba(0,0,0,0)",
        font=dict(family="Noto Sans KR, Malgun Gothic", color="#1F2937"),
    )
    fig1.update_traces(textposition="outside")
    st.plotly_chart(fig1, use_container_width=True)
    _excel_download(cust_sales, "고객사별_매출순위", "dl_cust_rank")

    st.subheader(f"고객사별 월별 매출 추이 (Top {top_n})")
    top_customers = (
        filtered.groupby("고객사명")[metric].sum().sort_values(ascending=False).head(top_n).index
    )
    cust_monthly = (
        filtered[filtered["고객사명"].isin(top_customers)]
        .groupby(["년월", "고객사명"])[metric]
        .sum()
        .reset_index()
    )
    fig2 = px.line(
        cust_monthly,
        x="년월",
        y=metric,
        color="고객사명",
        markers=True,
        color_discrete_sequence=PI_COLORS,
    )
    _pi_layout(fig2, height=450)
    st.plotly_chart(fig2, use_container_width=True)
    _excel_download(cust_monthly, "고객사별_월별추이", "dl_cust_monthly")

    # 고객사별 상세 테이블
    st.subheader(f"고객사별 매출 상세 (Top {top_n})")
    cust_detail = (
        filtered.groupby(["고객사명", "고객코드"])
        .agg({metric: "sum", "상품코드": "nunique"})
        .rename(columns={"상품코드": "상품수"})
        .sort_values(metric, ascending=False)
        .head(top_n)
        .reset_index()
    )
    cust_detail[metric] = cust_detail[metric].apply(lambda x: f"{x:,.0f}")
    st.dataframe(cust_detail, use_container_width=True, height=400)

# ============================================================
# 탭2: 품목(유형) 분석
# ============================================================
with tab2:
    col_a, col_b = st.columns(2)

    with col_a:
        st.subheader("중유형별 매출 순위")
        type_sales = (
            filtered.groupby("중유형")[metric]
            .sum()
            .sort_values(ascending=False)
            .reset_index()
        )
        fig3 = px.bar(
            type_sales,
            x="중유형",
            y=metric,
            text_auto=".3s",
            color=metric,
            color_continuous_scale=PI_GRADIENT,
        )
        _pi_layout(fig3, xaxis_tickangle=-45, coloraxis_showscale=False)
        fig3.update_traces(textposition="outside")
        st.plotly_chart(fig3, use_container_width=True)
        _excel_download(type_sales, "중유형별_매출", "dl_type")

    with col_b:
        st.subheader("소유형별 매출 순위")
        subtype_sales = (
            filtered.groupby("소유형")[metric]
            .sum()
            .sort_values(ascending=False)
            .head(15)
            .reset_index()
        )
        fig4 = px.bar(
            subtype_sales,
            x="소유형",
            y=metric,
            text_auto=".3s",
            color=metric,
            color_continuous_scale=[[0, "#9CA3AF"], [1, "#1F2937"]],
        )
        _pi_layout(fig4, xaxis_tickangle=-45, coloraxis_showscale=False)
        fig4.update_traces(textposition="outside")
        st.plotly_chart(fig4, use_container_width=True)
        _excel_download(subtype_sales, "소유형별_매출", "dl_subtype")

    st.subheader("중유형별 월별 매출 추이")
    type_monthly = (
        filtered.groupby(["년월", "중유형"])[metric].sum().reset_index()
    )
    fig5 = px.bar(
        type_monthly,
        x="년월",
        y=metric,
        color="중유형",
        barmode="group",
        text_auto=".3s",
        color_discrete_sequence=PI_COLORS,
    )
    _pi_layout(fig5, height=450)
    st.plotly_chart(fig5, use_container_width=True)
    _excel_download(type_monthly, "중유형별_월별추이", "dl_type_monthly")

    st.subheader("소유형별 월별 매출 추이")
    top_subtypes = (
        filtered.groupby("소유형")[metric].sum().sort_values(ascending=False).head(8).index
    )
    subtype_monthly = (
        filtered[filtered["소유형"].isin(top_subtypes)]
        .groupby(["년월", "소유형"])[metric]
        .sum()
        .reset_index()
    )
    fig5b = px.line(
        subtype_monthly,
        x="년월",
        y=metric,
        color="소유형",
        markers=True,
        color_discrete_sequence=PI_COLORS,
    )
    _pi_layout(fig5b, height=450)
    st.plotly_chart(fig5b, use_container_width=True)
    _excel_download(subtype_monthly, "소유형별_월별추이", "dl_subtype_monthly")

    # 19년 카테고리별
    st.subheader("19년 카테고리별 매출")
    cat_sales = (
        filtered.groupby("19년 카테고리")[metric]
        .sum()
        .sort_values(ascending=False)
        .reset_index()
    )
    fig_cat = px.pie(
        cat_sales,
        names="19년 카테고리",
        values=metric,
        hole=0.3,
        color_discrete_sequence=PI_COLORS,
    )
    _pi_layout(fig_cat, height=450)
    st.plotly_chart(fig_cat, use_container_width=True)
    _excel_download(cat_sales, "카테고리별_매출", "dl_cat")

# ============================================================
# 탭3: 국내/해외 분석
# ============================================================
with tab3:
    col_c, col_d = st.columns(2)

    with col_c:
        st.subheader("국내 vs 해외 매출 비중")
        region_sales = filtered.groupby("국내/해외")[metric].sum().reset_index()
        fig6 = px.pie(
            region_sales,
            names="국내/해외",
            values=metric,
            hole=0.4,
            color_discrete_map={"국내": "#1F2937", "해외": "#E61E3D"},
        )
        _pi_layout(fig6, height=400)
        st.plotly_chart(fig6, use_container_width=True)
        _excel_download(region_sales, "국내해외_비중", "dl_region_pie")

    with col_d:
        st.subheader("국내/해외 월별 매출 추이")
        region_monthly = (
            filtered.groupby(["년월", "국내/해외"])[metric].sum().reset_index()
        )
        fig7 = px.bar(
            region_monthly,
            x="년월",
            y=metric,
            color="국내/해외",
            barmode="group",
            text_auto=".3s",
            color_discrete_map={"국내": "#1F2937", "해외": "#E61E3D"},
        )
        _pi_layout(fig7, height=400)
        st.plotly_chart(fig7, use_container_width=True)
        _excel_download(region_monthly, "국내해외_월별추이", "dl_region_monthly")

    st.subheader(f"국내 - 고객사별 매출 Top {top_n}")
    dom = filtered[filtered["국내/해외"] == "국내"]
    dom_top = (
        dom.groupby("고객사명")[metric]
        .sum()
        .sort_values(ascending=False)
        .head(top_n)
        .reset_index()
    )
    fig8 = px.bar(
        dom_top, x="고객사명", y=metric, text_auto=".3s", color_discrete_sequence=["#1F2937"]
    )
    _pi_layout(fig8, xaxis_tickangle=-45)
    fig8.update_traces(textposition="outside")
    st.plotly_chart(fig8, use_container_width=True)
    _excel_download(dom_top, "국내_고객사_Top", "dl_dom_top")

    st.subheader(f"해외 - 고객사별 매출 Top {top_n}")
    ovs = filtered[filtered["국내/해외"] == "해외"]
    ovs_top = (
        ovs.groupby("고객사명")[metric]
        .sum()
        .sort_values(ascending=False)
        .head(top_n)
        .reset_index()
    )
    fig9 = px.bar(
        ovs_top, x="고객사명", y=metric, text_auto=".3s", color_discrete_sequence=["#E61E3D"]
    )
    _pi_layout(fig9, xaxis_tickangle=-45)
    fig9.update_traces(textposition="outside")
    st.plotly_chart(fig9, use_container_width=True)
    _excel_download(ovs_top, "해외_고객사_Top", "dl_ovs_top")

    # 국내/해외별 유형 비교
    st.subheader("국내/해외별 중유형 매출 비교")
    region_type = (
        filtered.groupby(["국내/해외", "중유형"])[metric].sum().reset_index()
    )
    fig10 = px.bar(
        region_type,
        x="중유형",
        y=metric,
        color="국내/해외",
        barmode="group",
        text_auto=".3s",
        color_discrete_map={"국내": "#1F2937", "해외": "#E61E3D"},
    )
    _pi_layout(fig10, xaxis_tickangle=-45, height=450)
    st.plotly_chart(fig10, use_container_width=True)
    _excel_download(region_type, "국내해외_중유형비교", "dl_region_type")

# ============================================================
# 탭5: 담당자 분석
# ============================================================
with tab5:
    st.subheader("담당자별 매출 순위 & 담당 제품수")

    mgr_summary = (
        filtered.groupby(["현담당자", "현담당자 팀"])
        .agg(
            매출합계=(metric, "sum"),
            담당제품수=("상품코드", "nunique"),
            담당고객수=("고객코드", "nunique"),
        )
        .sort_values("매출합계", ascending=False)
        .reset_index()
    )

    col_m1, col_m2 = st.columns(2)

    with col_m1:
        fig_mgr = px.bar(
            mgr_summary,
            x="현담당자",
            y="매출합계",
            color="현담당자 팀",
            text_auto=".3s",
            color_discrete_sequence=PI_COLORS,
        )
        _pi_layout(fig_mgr, xaxis_tickangle=-45, height=500)
        fig_mgr.update_traces(textposition="outside")
        st.plotly_chart(fig_mgr, use_container_width=True)
        _excel_download(mgr_summary[["현담당자", "현담당자 팀", "매출합계"]], "담당자별_매출", "dl_mgr_sales")

    with col_m2:
        fig_prod = px.bar(
            mgr_summary,
            x="현담당자",
            y="담당제품수",
            color="현담당자 팀",
            text_auto=True,
            color_discrete_sequence=PI_COLORS,
        )
        _pi_layout(fig_prod, xaxis_tickangle=-45, height=500)
        fig_prod.update_traces(textposition="outside")
        st.plotly_chart(fig_prod, use_container_width=True)
        _excel_download(mgr_summary[["현담당자", "현담당자 팀", "담당제품수", "담당고객수"]], "담당자별_제품수", "dl_mgr_prod")

    # 담당자별 상세 테이블
    st.subheader("담당자별 상세 현황")
    mgr_display = mgr_summary.copy()
    mgr_display["매출합계"] = mgr_display["매출합계"].apply(lambda x: f"{x:,.0f}")
    st.dataframe(mgr_display, use_container_width=True, height=400)

    # 담당자별 월별 추이
    st.subheader("담당자별 월별 매출 추이 (Top 10)")
    top_mgrs = (
        filtered.groupby("현담당자")[metric]
        .sum()
        .sort_values(ascending=False)
        .head(10)
        .index
    )
    mgr_monthly = (
        filtered[filtered["현담당자"].isin(top_mgrs)]
        .groupby(["년월", "현담당자"])[metric]
        .sum()
        .reset_index()
    )
    fig_mgr_trend = px.line(
        mgr_monthly,
        x="년월",
        y=metric,
        color="현담당자",
        markers=True,
        color_discrete_sequence=PI_COLORS,
    )
    _pi_layout(fig_mgr_trend, height=450)
    st.plotly_chart(fig_mgr_trend, use_container_width=True)
    _excel_download(mgr_monthly, "담당자별_월별추이", "dl_mgr_monthly")

    # 담당자별 고객사 구성
    st.subheader("담당자별 주요 고객사")
    selected_mgr = st.selectbox("담당자 선택", mgr_summary["현담당자"].tolist())
    if selected_mgr:
        mgr_cust = (
            filtered[filtered["현담당자"] == selected_mgr]
            .groupby("고객사명")
            .agg({metric: "sum", "상품코드": "nunique"})
            .rename(columns={"상품코드": "제품수"})
            .sort_values(metric, ascending=False)
            .reset_index()
        )
        fig_mgr_cust = px.bar(
            mgr_cust,
            x="고객사명",
            y=metric,
            text="제품수",
            color_discrete_sequence=["#E61E3D"],
        )
        fig_mgr_cust.update_traces(
            texttemplate="제품 %{text}건", textposition="outside"
        )
        _pi_layout(fig_mgr_cust, xaxis_tickangle=-45, height=400)
        st.plotly_chart(fig_mgr_cust, use_container_width=True)
        _excel_download(mgr_cust, f"담당자_{selected_mgr}_고객사", "dl_mgr_cust")

# ============================================================
# 탭7: 키워드 분석
# ============================================================
with tab7:

    # ── 키워드 사전 (효능/성분/제형 관련) ──
    KEYWORDS = [
        # 효능
        "미백", "화이트", "화이트닝", "브라이트", "브라이트닝", "톤업", "광채",
        "잡티", "다크스팟", "멜라닌",
        "주름", "링클", "탄력", "리프트", "리프팅", "안티에이징", "에이징",
        "보습", "수분", "하이드로", "하이드라", "모이스처", "히알루론",
        "진정", "시카", "카밍", "수딩", "센시티브",
        "각질", "필링", "스크럽",
        "모공", "포어", "세범", "피지",
        "자외선", "선", "UV", "SPF", "자차",
        # 성분
        "비타민", "비타", "비타씨", "레티놀", "나이아신", "콜라겐",
        "펩타이드", "펩타", "세라마이드",
        "프로폴리스", "꿀", "허니",
        "갈락토미", "발효", "프로바이오틱",
        "어성초", "티트리", "녹차",
        "달팽이", "뮤신",
        "병풀", "마데카", "센텔라",
        # 제형/타입
        "앰플", "세럼", "에센스", "토너", "스킨",
        "크림", "로션", "밀크", "젤",
        "마스크", "팩", "패드", "패치",
        "미스트", "스프레이",
        "클렌저", "클렌징", "폼", "워시",
        "오일", "밤", "버터",
    ]

    @st.cache_data
    def extract_keywords(data, metric_col):
        """상품명에서 키워드 추출 후 매출 집계"""
        records = []
        for _, row in data.iterrows():
            name = str(row["상품명"]).upper()
            found = set()
            for kw in KEYWORDS:
                if kw.upper() in name and kw.upper() not in found:
                    found.add(kw.upper())
                    records.append({
                        "키워드": kw,
                        "년도": row["년도"],
                        "월": row["월"],
                        "분기": f'Q{(int(row["월"])-1)//3+1}',
                        "국내/해외": row["국내/해외"],
                        metric_col: row[metric_col],
                        "상품코드": row["상품코드"],
                    })
        return pd.DataFrame(records) if records else pd.DataFrame()

    kw_df = extract_keywords(filtered, metric)

    if kw_df.empty:
        st.info("키워드를 포함하는 상품이 없습니다.")
    else:
        st.subheader("제품명 키워드 분석 (Top 10)")
        st.caption("상품명에서 효능/성분/제형 키워드를 자동 추출하여 매출을 집계합니다.")

        # ── 1) 전체 키워드 Top 10 ──
        kw_rank = kw_df.groupby("키워드").agg(
            매출합계=(metric, "sum"), 제품수=("상품코드", "nunique")
        ).sort_values("매출합계", ascending=False).head(10).reset_index()

        col_k1, col_k2 = st.columns(2)
        with col_k1:
            fig_kw = px.bar(
                kw_rank, x="키워드", y="매출합계",
                text_auto=".3s",
                color="매출합계",
                color_continuous_scale=PI_GRADIENT,
            )
            _pi_layout(fig_kw, height=450, coloraxis_showscale=False,
                       xaxis_tickangle=-45, title_text="키워드별 매출 Top 10")
            fig_kw.update_traces(textposition="outside")
            st.plotly_chart(fig_kw, use_container_width=True)
            _excel_download(kw_rank, "키워드_매출순위", "dl_kw_rank")

        with col_k2:
            fig_kw2 = px.bar(
                kw_rank, x="키워드", y="제품수",
                text_auto=True,
                color="제품수",
                color_continuous_scale=[[0, "#9CA3AF"], [1, "#1F2937"]],
            )
            _pi_layout(fig_kw2, height=450, coloraxis_showscale=False,
                       xaxis_tickangle=-45, title_text="키워드별 제품 수 Top 10")
            fig_kw2.update_traces(textposition="outside")
            st.plotly_chart(fig_kw2, use_container_width=True)

        st.divider()

        # ── 2) 월별 키워드 추이 ──
        st.subheader("키워드별 월별 매출 추이")
        top_kws = kw_rank["키워드"].tolist()
        kw_monthly = (
            kw_df[kw_df["키워드"].isin(top_kws)]
            .groupby(["년월" if "년월" in kw_df.columns else "월", "키워드"])[metric]
            .sum()
            .reset_index()
        )
        # 년월 컬럼 생성
        if "년월" not in kw_df.columns:
            kw_df["년월"] = kw_df["년도"].astype(str) + "-" + kw_df["월"].astype(str).str.zfill(2)
            kw_monthly = (
                kw_df[kw_df["키워드"].isin(top_kws)]
                .groupby(["년월", "키워드"])[metric]
                .sum()
                .reset_index()
            )

        fig_kw_trend = px.line(
            kw_monthly, x="년월", y=metric, color="키워드",
            markers=True, color_discrete_sequence=PI_COLORS,
        )
        _pi_layout(fig_kw_trend, height=450)
        st.plotly_chart(fig_kw_trend, use_container_width=True)
        _excel_download(kw_monthly, "키워드_월별추이", "dl_kw_monthly")

        st.divider()

        # ── 3) 분기별 키워드 히트맵 ──
        st.subheader("분기별 키워드 매출 히트맵")
        kw_quarter = (
            kw_df[kw_df["키워드"].isin(top_kws)]
            .groupby(["분기", "키워드"])[metric]
            .sum()
            .reset_index()
        )
        kw_q_pivot = kw_quarter.pivot(index="키워드", columns="분기", values=metric).fillna(0)
        # 매출 순 정렬
        kw_q_pivot = kw_q_pivot.loc[top_kws]
        kw_q_pivot = kw_q_pivot.reindex(columns=sorted(kw_q_pivot.columns))

        fig_kw_heat = go.Figure(data=go.Heatmap(
            z=kw_q_pivot.values,
            x=kw_q_pivot.columns.tolist(),
            y=kw_q_pivot.index.tolist(),
            colorscale=[[0, "#F9FAFB"], [0.3, "#9CA3AF"], [0.7, "#DC2626"], [1, "#E61E3D"]],
            hovertemplate="키워드: %{y}<br>분기: %{x}<br>매출: %{z:,.0f}<extra></extra>",
            showscale=True,
            colorbar=dict(title=dict(text="매출(원)", font=dict(size=10))),
        ))
        _pi_layout(fig_kw_heat, height=400,
                   yaxis=dict(autorange="reversed"),
                   xaxis_title="분기", yaxis_title="")
        st.plotly_chart(fig_kw_heat, use_container_width=True)
        _excel_download(kw_q_pivot.reset_index(), "키워드_분기별", "dl_kw_quarter")

        st.divider()

        # ── 4) 국내/해외 키워드 비교 ──
        st.subheader("국내 vs 해외 키워드 매출 비교")
        kw_region = (
            kw_df[kw_df["키워드"].isin(top_kws)]
            .groupby(["국내/해외", "키워드"])[metric]
            .sum()
            .reset_index()
        )
        fig_kw_region = px.bar(
            kw_region, x="키워드", y=metric,
            color="국내/해외", barmode="group",
            text_auto=".3s",
            color_discrete_map={"국내": "#1F2937", "해외": "#E61E3D"},
        )
        _pi_layout(fig_kw_region, height=450, xaxis_tickangle=-45)
        st.plotly_chart(fig_kw_region, use_container_width=True)
        _excel_download(kw_region, "키워드_국내해외", "dl_kw_region")

        st.divider()

        # ── 5) 키워드 워드클라우드 (테이블 형태) ──
        st.subheader("전체 키워드 매출 현황")
        kw_all = kw_df.groupby("키워드").agg(
            매출합계=(metric, "sum"), 제품수=("상품코드", "nunique")
        ).sort_values("매출합계", ascending=False).reset_index()
        kw_all["매출(억원)"] = (kw_all["매출합계"] / 1e8).round(1)
        total_kw = kw_all["매출합계"].sum()
        kw_all["비중(%)"] = (kw_all["매출합계"] / total_kw * 100).round(1) if total_kw else 0
        st.dataframe(
            kw_all[["키워드", "매출(억원)", "제품수", "비중(%)"]],
            use_container_width=True, height=400, hide_index=True,
        )
        _excel_download(kw_all, "키워드_전체현황", "dl_kw_all")

# ============================================================
# 탭6: 인터랙티브 분석
# ============================================================
with tab6:

    # ── 1) 트리맵: 매출 구조 한눈에 보기 ──
    st.subheader("매출 구조 트리맵")
    treemap_level = st.radio(
        "트리맵 계층 선택",
        ["국내/해외 → 중유형 → 고객사", "Team → 담당자 → 고객사", "중유형 → 소유형 → 고객사"],
        horizontal=True,
        key="treemap_level",
    )
    level_map = {
        "국내/해외 → 중유형 → 고객사": ["국내/해외", "중유형", "고객사명"],
        "Team → 담당자 → 고객사": ["Team", "현담당자", "고객사명"],
        "중유형 → 소유형 → 고객사": ["중유형", "소유형", "고객사명"],
    }
    tree_path = level_map[treemap_level]
    tree_data = filtered.groupby(tree_path)[metric].sum().reset_index()
    tree_data = tree_data[tree_data[metric] > 0].copy()
    fig_tree = px.treemap(
        tree_data,
        path=[px.Constant("전체")] + tree_path,
        values=metric,
        color=metric,
        color_continuous_scale=[[0, "#374151"], [0.4, "#6B7280"], [0.7, "#DC2626"], [1, "#E61E3D"]],
    )
    fig_tree.update_traces(
        textinfo="label+value+percent parent",
        texttemplate="<b>%{label}</b><br>%{value:,.0f}<br>(%{percentParent:.1%})",
    )
    _pi_layout(fig_tree, height=600, coloraxis_showscale=False)
    st.plotly_chart(fig_tree, use_container_width=True)
    _excel_download(tree_data, "트리맵_데이터", "dl_treemap")

    st.divider()

    # ── 2) 선버스트: 매출 비중 드릴다운 ──
    st.subheader("매출 비중 선버스트 (클릭하여 드릴다운)")
    sun_data = filtered.groupby(["국내/해외", "중유형", "소유형"])[metric].sum().reset_index()
    sun_data = sun_data[sun_data[metric] > 0].copy()
    fig_sun = px.sunburst(
        sun_data,
        path=["국내/해외", "중유형", "소유형"],
        values=metric,
        color=metric,
        color_continuous_scale=[[0, "#374151"], [0.5, "#9CA3AF"], [1, "#E61E3D"]],
    )
    fig_sun.update_traces(
        textinfo="label+percent parent",
        insidetextorientation="radial",
    )
    _pi_layout(fig_sun, height=550, coloraxis_showscale=False)
    st.plotly_chart(fig_sun, use_container_width=True)
    _excel_download(sun_data, "선버스트_데이터", "dl_sunburst")

    st.divider()

    # ── 3) 버블차트: 고객사별 매출 vs 제품수 ──
    st.subheader("고객사별 매출 vs 제품수 (버블 크기 = 매출)")
    bubble_data = (
        filtered.groupby(["고객사명", "국내/해외"])
        .agg(매출합계=(metric, "sum"), 제품수=("상품코드", "nunique"), 담당자수=("현담당자", "nunique"))
        .reset_index()
    )
    bubble_data = bubble_data[bubble_data["매출합계"] > 0].copy()
    bubble_data["제품당매출"] = (bubble_data["매출합계"] / bubble_data["제품수"].replace(0, 1))
    fig_bubble = px.scatter(
        bubble_data,
        x="제품수",
        y="매출합계",
        size="매출합계",
        color="국내/해외",
        hover_name="고객사명",
        hover_data={"제품당매출": ":,.0f", "담당자수": True},
        color_discrete_map={"국내": "#1F2937", "해외": "#E61E3D"},
        size_max=60,
    )
    _pi_layout(fig_bubble, height=550,
               xaxis_title="제품 수",
               yaxis_title=f"{metric} (원)")
    st.plotly_chart(fig_bubble, use_container_width=True)
    _excel_download(bubble_data, "버블차트_데이터", "dl_bubble")

    st.divider()

    # ── 4) 히트맵: 고객사 × 월 매출 ──
    st.subheader("고객사 × 월 매출 히트맵 (클릭/줌 가능)")
    heatmap_n = st.slider("히트맵 표시 고객사 수", 5, 30, 15, key="heatmap_n")
    top_heat_cust = filtered.groupby("고객사명")[metric].sum().sort_values(ascending=False).head(heatmap_n).index
    heat_data = (
        filtered[filtered["고객사명"].isin(top_heat_cust)]
        .groupby(["고객사명", "년월"])[metric]
        .sum()
        .reset_index()
    )
    heat_pivot = heat_data.pivot(index="고객사명", columns="년월", values=metric).fillna(0)
    # 매출 순 정렬
    heat_pivot = heat_pivot.loc[top_heat_cust]

    fig_heat = go.Figure(data=go.Heatmap(
        z=heat_pivot.values,
        x=heat_pivot.columns.tolist(),
        y=heat_pivot.index.tolist(),
        colorscale=[[0, "#F9FAFB"], [0.3, "#9CA3AF"], [0.6, "#DC2626"], [1, "#E61E3D"]],
        hovertemplate="고객사: %{y}<br>월: %{x}<br>매출: %{z:,.0f}<extra></extra>",
        showscale=True,
        colorbar=dict(title=dict(text="매출(원)", font=dict(size=10))),
    ))
    _pi_layout(fig_heat, height=max(400, heatmap_n * 28),
               yaxis=dict(autorange="reversed"),
               xaxis_title="월", yaxis_title="")
    st.plotly_chart(fig_heat, use_container_width=True)
    _excel_download(heat_pivot.reset_index(), "히트맵_데이터", "dl_heatmap")

    st.divider()

    # ── 5) 산키 다이어그램: 매출 흐름 ──
    st.subheader("매출 흐름도 (국내/해외 → 중유형 → Top 고객사)")
    sankey_top = 8
    top_sankey_cust = filtered.groupby("고객사명")[metric].sum().sort_values(ascending=False).head(sankey_top).index

    sankey_df = filtered.copy()
    sankey_df["고객사표시"] = sankey_df["고객사명"].apply(lambda x: x if x in top_sankey_cust else "기타")

    # 1단계: 국내/해외 → 중유형
    flow1 = sankey_df.groupby(["국내/해외", "중유형"])[metric].sum().reset_index()
    flow1 = flow1[flow1[metric] > 0]
    # 2단계: 중유형 → 고객사
    flow2 = sankey_df.groupby(["중유형", "고객사표시"])[metric].sum().reset_index()
    flow2 = flow2[flow2[metric] > 0]

    # 노드 목록 생성
    nodes_list = []
    for col in ["국내/해외", "중유형", "고객사표시"]:
        for val in sankey_df[col].unique():
            if val not in nodes_list:
                nodes_list.append(val)
    node_idx = {n: i for i, n in enumerate(nodes_list)}

    sources, targets, values = [], [], []
    for _, r in flow1.iterrows():
        sources.append(node_idx[r["국내/해외"]])
        targets.append(node_idx[r["중유형"]])
        values.append(r[metric])
    for _, r in flow2.iterrows():
        sources.append(node_idx[r["중유형"]])
        targets.append(node_idx[r["고객사표시"]])
        values.append(r[metric])

    # 노드 색상
    node_colors = []
    for n in nodes_list:
        if n == "국내":
            node_colors.append("#1F2937")
        elif n == "해외":
            node_colors.append("#E61E3D")
        elif n in sankey_df["중유형"].unique():
            node_colors.append("#6B7280")
        else:
            node_colors.append("#DC2626")

    fig_sankey = go.Figure(data=[go.Sankey(
        textfont=dict(size=13, color="#111827", family="Noto Sans KR, Malgun Gothic"),
        node=dict(
            pad=20, thickness=25,
            label=nodes_list,
            color=node_colors,
        ),
        link=dict(
            source=sources,
            target=targets,
            value=values,
            color="rgba(180,180,180,0.4)",
        ),
    )])
    _pi_layout(fig_sankey, height=550, title_text="")
    st.plotly_chart(fig_sankey, use_container_width=True)

# ============================================================
# 탭4: 상세 데이터
# ============================================================
with tab4:
    st.subheader("원본 데이터 조회")

    # 고객사 검색
    search = st.text_input("고객사명 또는 상품명 검색")
    display_df = filtered.copy()
    if search:
        mask = (
            display_df["고객사명"].str.contains(search, case=False, na=False)
            | display_df["상품명"].str.contains(search, case=False, na=False)
            | display_df["상품코드"].str.contains(search, case=False, na=False)
        )
        display_df = display_df[mask]

    display_cols = [
        "국내/해외", "년월", "고객사명", "상품코드", "상품명",
        "소유형", "중유형", "제품매출", "순매출액", "Team", "현담당자",
    ]
    st.dataframe(
        display_df[display_cols].sort_values(metric, ascending=False),
        use_container_width=True,
        height=600,
    )

    st.download_button(
        "📥 필터링 데이터 다운로드 (CSV)",
        display_df.to_csv(index=False).encode("utf-8-sig"),
        "filtered_data.csv",
        "text/csv",
    )

# ============================================================
# 탭: 작성 방법
# ============================================================
with tab_help:
    st.subheader("엑셀 파일 작성 방법")

    st.markdown("""
이 대시보드는 **엑셀 파일(.xlsx)**을 업로드하면 자동으로 매출 분석을 수행합니다.
아래 규칙에 맞게 엑셀 파일을 작성해주세요.
    """)

    st.markdown("---")

    st.markdown("### 1. 시트 구성")
    st.markdown("""
| 시트명 | 용도 | 필수 여부 |
|---|---|---|
| **ES1** | ES1팀 매출 데이터 | 필수 (최소 1개 시트) |
| **ES2** | ES2팀 매출 데이터 | 선택 |
| **ES1 요약**, **ES2 요약** | 요약 피벗 (자동 무시됨) | 불필요 |

- 시트명에 **"요약"**이 포함된 시트는 자동으로 제외됩니다.
- 시트 수는 자유롭게 추가 가능합니다 (ES3, ES4 등).
    """)

    st.markdown("### 2. 필수 컬럼 (21개)")

    col_info = pd.DataFrame([
        ["국내/해외", "국내 / 해외", "국내", "필수"],
        ["년도", "연도 (숫자)", "2026", "필수"],
        ["월", "월 (숫자)", "1", "필수"],
        ["상품코드", "상품 고유코드 (고객코드 추출용)", "9ABC1180815", "필수"],
        ["제품 관리유형 코드", "관리유형 코드", "11S0701", "선택"],
        ["상품명", "상품 이름", "미샤비타씨플러스앰플", "필수"],
        ["벌크코드", "벌크 코드", "3ABC11808150", "선택"],
        ["벌크명", "벌크 이름", "미샤비타씨플러스앰플벌크", "선택"],
        ["출시연도", "출시 연도", "2025", "선택"],
        ["출시월", "출시 월", "6", "선택"],
        ["벌크관리유형코드", "벌크 관리유형", "13S0701", "선택"],
        ["현담당자", "담당 연구원", "김연구", "필수"],
        ["현담당자 팀", "담당자 소속 팀", "ES1팀", "필수"],
        ["제품매출", "제품매출액 (원)", "17694833", "필수"],
        ["순매출액", "순매출액 (원)", "17694833", "필수"],
        ["소유형", "제품 소분류", "에센스", "필수"],
        ["중유형", "제품 중분류", "에센스", "필수"],
        ["R&I", "R&I 구분", "S", "선택"],
        ["19년 카테고리", "카테고리 분류", "스킨/에센스류", "선택"],
        ["Lab", "연구소", "ES Lab", "선택"],
        ["Team", "팀 (ES1/ES2 등)", "ES1", "필수"],
    ], columns=["컬럼명", "설명", "예시", "필수여부"])
    st.dataframe(col_info, use_container_width=True, height=500, hide_index=True)

    st.markdown("### 3. 상품코드 규칙 (고객사 자동 추출)")
    st.markdown("""
상품코드의 **2~4번째 글자**가 고객사 코드로 자동 추출됩니다.

```
상품코드: 9ABC1180815
           ^^^
           고객코드 = ABC → 미샤(에이블씨엔씨)
```

| 상품코드 | 고객코드 | 고객사명 |
|---|---|---|
| 9**ABC**1180815 | ABC | 미샤(에이블씨엔씨) |
| 9**DPD**0001210 | DPD | 아누아 |
| 9**CLO**0012310 | CLO | 클리오/구달 |
| 9**GDI**0005610 | GDI | 조선미녀 |

주요 고객사 약 60개는 자동 매핑되며, 매핑되지 않은 코드는 코드 그대로 표시됩니다.
    """)

    st.markdown("### 4. 주의사항")
    st.markdown("""
- **첫 번째 행은 헤더**(컬럼명)여야 합니다.
- **매출 금액은 숫자**(원 단위)로 입력하세요. 문자가 섞이면 0으로 처리됩니다.
- **음수 매출**(반품 등)도 허용되지만, 일부 차트(트리맵, 버블)에서는 제외됩니다.
- 여러 시트의 데이터는 자동으로 **합산**됩니다.
    """)

    st.markdown("### 5. 예제 파일 다운로드")
    st.markdown("아래 버튼으로 예제 파일을 다운로드하여 형식을 확인할 수 있습니다.")

    sample_path = os.path.join(os.path.dirname(__file__), "예제_매출데이터_v3.xlsx")
    if os.path.exists(sample_path):
        with open(sample_path, "rb") as f:
            st.download_button(
                "📥 예제 엑셀 파일 다운로드",
                data=f.read(),
                file_name="예제_매출데이터_v3.xlsx",
                mime="application/vnd.openxmlformats-officedocument.spreadsheetml.sheet",
                key="dl_sample",
            )

    st.markdown("### 6. 대시보드 기능 안내")
    st.markdown("""
| 탭 | 기능 |
|---|---|
| **📊 실적 현황표** | 월별 비교 분석표 (▲▼ 증감 표시), 고객사/유형/담당자/Team 기준 전환 |
| **🏢 고객사 분석** | 고객사별 매출 순위, 월별 추이, 상세 테이블 |
| **📦 품목 분석** | 중유형/소유형별 순위, 월별 추이, 카테고리 파이차트 |
| **🌏 국내/해외** | 국내 vs 해외 비중, 각각 Top N 고객사, 유형 비교 |
| **👤 담당자 분석** | 담당자별 매출/제품수, 월별 추이, 고객사 구성 |
| **🔍 인터랙티브** | 트리맵, 선버스트, 버블차트, 히트맵, 산키 다이어그램 |
| **📋 상세 데이터** | 원본 데이터 검색/필터링 + CSV 다운로드 |

**사이드바 필터**: 국내/해외, 월, Team, 매출기준(제품매출/순매출액), 상위 고객사 수

모든 차트 아래 **📥 엑셀 다운로드** 버튼으로 해당 데이터를 받을 수 있습니다.
    """)
